import io

from pptx import Presentation

from app.ingest.extractors.chunking import (
    DEFAULT_CHUNK_OVERLAP,
    DEFAULT_CHUNK_SIZE,
    chunk_text,
)


def _extract_slide_text(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is not None:
            texts.append(notes_frame.text)
    return texts


class PptxChunkExtractor:
    """Extract text from PPTX bytes and chunk for embedding."""

    def __init__(
        self,
        chunk_size: int = DEFAULT_CHUNK_SIZE,
        chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
    ) -> None:
        self._chunk_size = chunk_size
        self._chunk_overlap = chunk_overlap

    def extract(self, data: bytes, source: str) -> list[dict]:
        prs = Presentation(io.BytesIO(data))
        parts = [
            "\n".join(slide_texts)
            for slide in prs.slides
            if (slide_texts := _extract_slide_text(slide))
        ]
        full_text = "\n\n".join(parts)
        chunks = chunk_text(
            full_text,
            chunk_size=self._chunk_size,
            chunk_overlap=self._chunk_overlap,
        )
        return [{"text": chunk, "source": source} for chunk in chunks]
